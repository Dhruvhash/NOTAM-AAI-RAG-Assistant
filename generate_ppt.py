import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # Set 16:9 widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Theme Colors
    NAVY = RGBColor(15, 23, 42)         # #0F172A
    DARK_BLUE = RGBColor(30, 58, 138)   # #1E3A8A
    PRIMARY_BLUE = RGBColor(37, 99, 235) # #2563EB
    CYAN_ACCENT = RGBColor(14, 165, 233)# #0EA5E9
    BG_LIGHT = RGBColor(248, 250, 252)  # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)   # #FFFFFF
    CARD_BORDER = RGBColor(226, 232, 240)# #E2E8F0
    TEXT_MAIN = RGBColor(30, 41, 59)    # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)# #64748B
    WHITE = RGBColor(255, 255, 255)
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981
    ACCENT_ORANGE = RGBColor(245, 158, 11) # #F59E0B

    def add_header(slide, title_text, category_text="AIRPORTS AUTHORITY OF INDIA | IT DEPARTMENT"):
        # Top header accent line
        header_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = CYAN_ACCENT
        header_line.line.color.rgb = CYAN_ACCENT

        # Category pill text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        p.font.name = "Calibri"

        # Main slide title
        txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        p_title.font.name = "Calibri"

    def add_footer(slide, current_page, total_pages=6):
        # Footer text box
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = "AI-Powered NOTAM Assistant | Retrieval Augmented Generation (RAG) System"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.font.name = "Calibri"
        
        # Right aligned page number
        p_num = tf.add_paragraph()
        p_num.text = f"Slide {current_page} of {total_pages}"
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.font.size = Pt(9)
        p_num.font.bold = True
        p_num.font.color.rgb = PRIMARY_BLUE
        p_num.font.name = "Calibri"

    # =========================================================================
    # SLIDE 1: Title Slide (Dark Theme for Executive Impact)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Top accent bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = CYAN_ACCENT
    top_bar.line.fill.background()

    # Organization Badge Card
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(5.8), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.color.rgb = CYAN_ACCENT
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "AIRPORTS AUTHORITY OF INDIA (AAI) | NORTHERN REGION IT DEPT"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = CYAN_ACCENT
    p_b.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(1.8))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "AI-Powered NOTAM Assistant"
    p_t.font.size = Pt(44)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.font.name = "Calibri"

    p_sub = tf_t.add_paragraph()
    p_sub.text = "A Retrieval Augmented Generation (RAG) System for Aviation Notice to Airmen Decoding & Query Resolution"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = CYAN_ACCENT
    p_sub.font.name = "Calibri"

    # Divider line
    div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.0), Inches(11.333), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = CYAN_ACCENT
    div.line.fill.background()

    # Highlight Cards (3 key summary points)
    cards_data = [
        ("Problem Addressed", "Dense, abbreviated NOTAM text blocks cause high reading overhead and delay critical pre-flight interpretation."),
        ("Technical Innovation", "Combines PyMuPDF, SentenceTransformers, ChromaDB vector store, and Groq Llama 3.1 LLM."),
        ("Core Guarantee", "Strict source grounding to eliminate AI hallucination and provide verifiable, source-cited responses.")
    ]
    for i, (title, desc) in enumerate(cards_data):
        left_pos = Inches(1.0 + i * 3.9)
        c_shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(4.3), Inches(3.6), Inches(2.3))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate Dark
        c_shape.line.color.rgb = PRIMARY_BLUE
        
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.2)
        
        p1 = tf_c.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = CYAN_ACCENT
        
        p2 = tf_c.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(226, 232, 240)

    # =========================================================================
    # SLIDE 2: Problem & Operational Context
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = BG_LIGHT
    bg2.line.fill.background()
    add_header(slide2, "Background & Operational Problem Statement")
    add_footer(slide2, 2)

    # Left Column: Operational Challenge
    card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = CARD_BORDER
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = Inches(0.3)

    p_lh = tf_l.paragraphs[0]
    p_lh.text = "The NOTAM Information Bottleneck"
    p_lh.font.size = Pt(18)
    p_lh.font.bold = True
    p_lh.font.color.rgb = DARK_BLUE

    points_l = [
        ("Dense Shorthand Standards", "NOTAMs follow ICAO formats using complex Q-codes, contractions (e.g., TWY C CLSD, ILS GP OTS UFN), and raw geographic coordinates."),
        ("High Cognitive Burden", "Aviation personnel (pilots, dispatchers, ATC) must manually parse multi-page PDF bulletins under time-sensitive pre-flight pressure."),
        ("Search & Retrieval Deficit", "Traditional static PDF bulletins lack semantic search capabilities, forcing staff to scan entire text blocks to locate relevant runway or airspace restrictions."),
        ("Risk of Misinterpretation", "Manual decoding increases likelihood of missing critical hazard updates or misinterpreting validity time windows across multiple notices.")
    ]
    for title, text in points_l:
        p_t = tf_l.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = TEXT_MAIN
        
        # Append text inline
        run = p_t.add_run()
        run.text = text
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED

    # Right Column: The AI Imperative & Safety Boundaries
    card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = CARD_BORDER
    tf_r = card_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = Inches(0.3)

    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Safety-First RAG Approach vs General AI"
    p_rh.font.size = Pt(18)
    p_rh.font.bold = True
    p_rh.font.color.rgb = DARK_BLUE

    points_r = [
        ("The Hallucination Risk", "General-purpose conversational AI models can generate plausible-sounding but completely fabricated operational answers—unacceptable in aviation safety."),
        ("Retrieval Augmented Generation (RAG)", "Restricts language model generation exclusively to retrieved, verified NOTAM bulletin records stored in a vector database."),
        ("100% Traceability & Citation", "Every generated plain-English answer is explicitly linked to the exact source NOTAM chunk and raw text for human verification."),
        ("Explicit Fallback Mechanism", "If information is absent from bulletins, the assistant explicitly states 'No relevant NOTAM found' rather than guessing.")
    ]
    for title, text in points_r:
        p_t = tf_r.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = TEXT_MAIN
        
        run = p_t.add_run()
        run.text = text
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: System Architecture (3-Tier Service Breakdown)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = BG_LIGHT
    bg3.line.fill.background()
    add_header(slide3, "System Architecture & Microservice Technology Stack")
    add_footer(slide3, 3)

    tiers = [
        ("1. Presentation Tier (Frontend)", "React 19 & Tailwind CSS", PRIMARY_BLUE, [
            ("Interactive UI", "Dashboard overview, Live NOTAM feed, Chat view, Bookmarks"),
            ("Theme System", "Light & Dark modes optimized for low-light flight deck visibility"),
            ("Role-Based Views", "Tailored navigation for Pilot, Controller, Dispatcher, and Admin"),
            ("Asynchronous UI", "Real-time job progress tracking during multi-page PDF processing")
        ]),
        ("2. Application Gateway Tier", "Node.js & Express API Gateway", DARK_BLUE, [
            ("JWT Authentication", "HttpOnly cookies & bcrypt password hashing for secure sessions"),
            ("Structured Data Store", "MongoDB + Mongoose for users, metadata & saved bookmarks"),
            ("Service Proxy", "Routes uploads & natural language queries to Python AI backend"),
            ("Decoupled Design", "Ensures interface remains functional even if AI service reboots")
        ]),
        ("3. AI & RAG Microservice Tier", "Python & FastAPI Ecosystem", NAVY, [
            ("PDF Parser & Ingestion", "PyMuPDF text extraction with regex pattern-based segmentation"),
            ("Contraction Decoder", "Batch LLM decoding of dense ICAO shorthand to plain text"),
            ("Vector Storage Engine", "SentenceTransformers (all-MiniLM-L6-v2) + ChromaDB (384-dim)"),
            ("Grounded Generation", "Groq Cloud (Llama 3.1) executing source-constrained prompts")
        ])
    ]

    for i, (tier_title, tech_stack, color, items) in enumerate(tiers):
        left_pos = Inches(0.8 + i * 3.95)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Header banner inside card
        banner = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.1), Inches(1.6), Inches(3.55), Inches(0.75))
        banner.fill.solid()
        banner.fill.fore_color.rgb = color
        banner.line.fill.background()
        tf_bn = banner.text_frame
        tf_bn.word_wrap = True
        p_bn1 = tf_bn.paragraphs[0]
        p_bn1.text = tier_title
        p_bn1.font.size = Pt(13)
        p_bn1.font.bold = True
        p_bn1.font.color.rgb = WHITE

        p_bn2 = tf_bn.add_paragraph()
        p_bn2.text = tech_stack
        p_bn2.font.size = Pt(11)
        p_bn2.font.color.rgb = CYAN_ACCENT

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0.9)
        tf_c.margin_left = tf_c.margin_right = Inches(0.2)

        for item_title, item_desc in items:
            p_it = tf_c.add_paragraph()
            p_it.text = f"• {item_title}"
            p_it.font.bold = True
            p_it.font.size = Pt(11)
            p_it.font.color.rgb = TEXT_MAIN

            p_id = tf_c.add_paragraph()
            p_id.text = f"  {item_desc}"
            p_id.font.size = Pt(10)
            p_id.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: End-to-End RAG Data Workflow & Core Modules
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = BG_LIGHT
    bg4.line.fill.background()
    add_header(slide4, "End-to-End RAG Processing Workflow & Functional Modules")
    add_footer(slide4, 4)

    # Workflow Steps Bar (4 Horizontal Steps)
    steps = [
        ("1. Ingestion & Split", "PyMuPDF extracts PDF text; pattern matching segments bulletin into discrete NOTAM records."),
        ("2. Batch Decoding", "LLM converts dense ICAO contractions to structured plain English in high-throughput batches."),
        ("3. Vector Indexing", "SentenceTransformers generates 384-dim embeddings stored with metadata in ChromaDB."),
        ("4. Grounded Q&A", "Semantic query match retrieves top context; Groq Llama 3.1 generates source-cited answer.")
    ]

    for i, (s_title, s_desc) in enumerate(steps):
        left_pos = Inches(0.8 + i * 2.95)
        step_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(2.75), Inches(1.7))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else PRIMARY_BLUE
        step_box.line.color.rgb = CYAN_ACCENT

        tf_s = step_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = Inches(0.15)
        
        p_st = tf_s.paragraphs[0]
        p_st.text = s_title
        p_st.font.size = Pt(13)
        p_st.font.bold = True
        p_st.font.color.rgb = WHITE

        p_sd = tf_s.add_paragraph()
        p_sd.text = s_desc
        p_sd.font.size = Pt(10)
        p_sd.font.color.rgb = RGBColor(241, 245, 249)

    # Detailed Modules Matrix Card
    mod_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(11.733), Inches(3.3))
    mod_card.fill.solid()
    mod_card.fill.fore_color.rgb = CARD_BG
    mod_card.line.color.rgb = CARD_BORDER

    tf_m = mod_card.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = tf_m.margin_right = tf_m.margin_top = Inches(0.25)

    p_mh = tf_m.paragraphs[0]
    p_mh.text = "Key System Functional Modules"
    p_mh.font.size = Pt(16)
    p_mh.font.bold = True
    p_mh.font.color.rgb = DARK_BLUE

    modules_data = [
        ("PDF Bulletin Processing", "Handles multi-page PDF uploads, handles minor layout variances, extracts structured NOTAM blocks."),
        ("Contraction Decoder", "Expands standardized aviation contractions while retaining exact time windows, identifiers, and coordinates."),
        ("Hybrid Metadata Filtering", "Combines ChromaDB vector semantic similarity with MongoDB metadata filtering (airport code, severity)."),
        ("Analytics Dashboard", "Provides visual breakdown of active NOTAMs by severity (Critical, Warning, Info) and airport location."),
        ("User Bookmarking", "Enables pilots and dispatchers to persist specific decoded NOTAMs for instant pre-flight briefing retrieval."),
        ("RBAC & Security Controls", "Implements role-based access for Pilots, Controllers, Dispatchers, and Admins with HTTP-only cookies.")
    ]

    # Render as 2 columns of 3 items
    for idx, (m_name, m_desc) in enumerate(modules_data):
        col = idx // 3
        row = idx % 3
        
        # Add text box inside the mod_card layout
        tb_mod = slide4.shapes.add_textbox(Inches(1.1 + col * 5.7), Inches(4.0 + row * 0.85), Inches(5.4), Inches(0.75))
        tf_tb = tb_mod.text_frame
        tf_tb.word_wrap = True
        tf_tb.margin_top = tf_tb.margin_left = tf_tb.margin_right = tf_tb.margin_bottom = 0
        
        p = tf_tb.paragraphs[0]
        p.text = f"✔ {m_name}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = PRIMARY_BLUE
        
        r = p.add_run()
        r.text = m_desc
        r.font.bold = False
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 5: Testing Results & Key Engineering Observations
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = BG_LIGHT
    bg5.line.fill.background()
    add_header(slide5, "System Testing, Results & Engineering Observations")
    add_footer(slide5, 5)

    # 3 Summary Result Cards
    res_cards = [
        ("Decoding Fidelity", ACCENT_GREEN, "Sample NOTAM Decoding", [
            ("Raw Input", "TWY C CLSD FOR MAINT WEF 2407220400 TO 2407221200."),
            ("Decoded Output", "Taxiway C is closed for maintenance from 22 July 2024 at 04:00 UTC to 12:00 UTC."),
            ("Fidelity Verification", "Preserved exact dates, times, and taxiway identifier without altering operational context.")
        ]),
        ("Grounded Accuracy", PRIMARY_BLUE, "RAG Chat Verification", [
            ("Direct Operational Query", "'Is Runway 09 closed at VIDP?' -> Correctly retrieved relevant notice & cited source ID."),
            ("Out-of-Scope Query", "'Provide weather forecast' -> Refused answer & stated info not present in indexed NOTAMs."),
            ("Safety Significance", "Zero hallucination observed under strict system prompt grounding constraints.")
        ]),
        ("Performance Optimization", ACCENT_ORANGE, "Latency & Microservices", [
            ("Batch Decoding", "Batching NOTAM blocks into single LLM inference calls dramatically reduced document ingestion time."),
            ("Local Embeddings", "Local SentenceTransformers execution removed dependency on external embedding APIs."),
            ("Service Decoupling", "Independent testing of FastAPI, Express Gateway, and React UI streamlined bug isolation.")
        ])
    ]

    for i, (c_title, c_accent, sub_header, items) in enumerate(res_cards):
        left_pos = Inches(0.8 + i * 3.95)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Header bar
        hbar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.5), Inches(3.75), Inches(0.1))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = c_accent
        hbar.line.fill.background()

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0.25)
        tf_c.margin_left = tf_c.margin_right = Inches(0.2)

        p_h = tf_c.paragraphs[0]
        p_h.text = c_title
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = DARK_BLUE

        p_sub = tf_c.add_paragraph()
        p_sub.text = sub_header.upper()
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = c_accent

        for label, desc in items:
            p_l = tf_c.add_paragraph()
            p_l.text = f"• {label}: "
            p_l.font.bold = True
            p_l.font.size = Pt(11)
            p_l.font.color.rgb = TEXT_MAIN

            run = p_l.add_run()
            run.text = desc
            run.font.bold = False
            run.font.size = Pt(10.5)
            run.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: Conclusion & Future Scope Roadmap
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = NAVY
    bg6.line.fill.background()

    # Top accent bar
    top_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar6.fill.solid()
    top_bar6.fill.fore_color.rgb = CYAN_ACCENT
    top_bar6.line.fill.background()

    # Title box
    txBox6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.9))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    p6_cat = tf6.paragraphs[0]
    p6_cat.text = "SUMMARY & STRATEGIC ROADMAP".upper()
    p6_cat.font.size = Pt(10)
    p6_cat.font.bold = True
    p6_cat.font.color.rgb = CYAN_ACCENT

    p6_t = tf6.add_paragraph()
    p6_t.text = "Conclusion & Future Development Roadmap"
    p6_t.font.size = Pt(26)
    p6_t.font.bold = True
    p6_t.font.color.rgb = WHITE

    # Left Card: Summary of Internship Achievements
    card_achieve = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1))
    card_achieve.fill.solid()
    card_achieve.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card_achieve.line.color.rgb = PRIMARY_BLUE
    tf_ac = card_achieve.text_frame
    tf_ac.word_wrap = True
    tf_ac.margin_left = tf_ac.margin_right = tf_ac.margin_top = Inches(0.25)

    p_ach = tf_ac.paragraphs[0]
    p_ach.text = "Key Internship Deliverables & Outcomes"
    p_ach.font.size = Pt(16)
    p_ach.font.bold = True
    p_ach.font.color.rgb = CYAN_ACCENT

    achievements = [
        ("Working End-to-End Prototype", "Successfully built and demonstrated a full-stack RAG assistant for NOTAM bulletin processing during 5-week training."),
        ("Source-Grounded Accuracy", "Established strict prompt boundaries preventing LLM hallucinations and ensuring complete answer traceability."),
        ("Production-Style Microservices", "Engineered a scalable 3-tier architecture (React, Node.js, Python/FastAPI) separating presentation, routing, and AI logic."),
        ("Domain-Oriented Learning", "Gained practical experience bridging CS engineering concepts with real-world public-sector aviation requirements at AAI NR Headquarters.")
    ]
    for title, desc in achievements:
        p = tf_ac.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        r = p.add_run()
        r.text = desc
        r.font.bold = False
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(203, 213, 225)

    # Right Box: Phased Future Roadmap (3 Phases)
    card_road = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1))
    card_road.fill.solid()
    card_road.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card_road.line.color.rgb = CYAN_ACCENT
    tf_rd = card_road.text_frame
    tf_rd.word_wrap = True
    tf_rd.margin_left = tf_rd.margin_right = tf_rd.margin_top = Inches(0.25)

    p_rdh = tf_rd.paragraphs[0]
    p_rdh.text = "Phased Operational Roadmap"
    p_rdh.font.size = Pt(16)
    p_rdh.font.bold = True
    p_rdh.font.color.rgb = CYAN_ACCENT

    phases = [
        ("Phase 1: Validation & Hardening (Short-Term)", [
            "Formal accuracy evaluation of contraction decoding against larger AAI benchmark datasets.",
            "Enterprise security review, penetration testing, and identity integration."
        ]),
        ("Phase 2: Live Integration (Medium-Term)", [
            "Integration with live operational NOTAM feeds (FAA/ICAO Digital NOTAM APIs).",
            "Transition from local ChromaDB/MongoDB to managed cloud cluster infrastructure."
        ]),
        ("Phase 3: Ecosystem Expansion (Long-Term)", [
            "Broadening RAG scope to AIP supplements, circulars, and weather advisories.",
            "Developing native mobile applications for pilot cockpits and handheld dispatch units."
        ])
    ]
    for ph_title, ph_points in phases:
        p_ph = tf_rd.add_paragraph()
        p_ph.text = ph_title
        p_ph.font.bold = True
        p_ph.font.size = Pt(12)
        p_ph.font.color.rgb = PRIMARY_BLUE

        for pt in ph_points:
            p_pt = tf_rd.add_paragraph()
            p_pt.text = f"  • {pt}"
            p_pt.font.size = Pt(10.5)
            p_pt.font.color.rgb = RGBColor(226, 232, 240)

    # Footer for slide 6
    txBox6_ft = slide6.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.3))
    tf6_ft = txBox6_ft.text_frame
    p_ft = tf6_ft.paragraphs[0]
    p_ft.text = "AI-Powered NOTAM Assistant | Airports Authority of India (AAI) Internship Project"
    p_ft.font.size = Pt(9)
    p_ft.font.color.rgb = RGBColor(148, 163, 184)
    p_ft_num = tf6_ft.add_paragraph()
    p_ft_num.text = "Slide 6 of 6"
    p_ft_num.alignment = PP_ALIGN.RIGHT
    p_ft_num.font.size = Pt(9)
    p_ft_num.font.bold = True
    p_ft_num.font.color.rgb = CYAN_ACCENT

    # Save presentation
    output_filename = "AAI_NOTAM_Assistant_6_Slide_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved to {output_filename}")

if __name__ == "__main__":
    create_deck()
